import os
import io
import re
import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from youtube_transcript_api import YouTubeTranscriptApi
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.chat_models import ChatOpenAI
from langchain_community.vectorstores import Qdrant
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import qdrant_client

load_dotenv()

# Set the OpenAI API key
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
QDRANT_URL = os.getenv('QDRANT_URL')
QDRANT_API_KEY = os.getenv('QDRANT_API_KEY')

# Initialize Qdrant Client
qdrant_client_instance = qdrant_client.QdrantClient(
    url=QDRANT_URL, api_key=QDRANT_API_KEY
)

# Initialize Qdrant Vector Store
def get_qdrant_vector_store():
    embeddings = OpenAIEmbeddings(
        model="text-embedding-ada-002", openai_api_key=OPENAI_API_KEY
    )
    return Qdrant(
        client=qdrant_client_instance,
        collection_name="vector-db-streamlit",
        embeddings=embeddings,
    )

vector_store = get_qdrant_vector_store()

# Text extraction functions
def get_pdf_text(pdf_bytes):
    text = ""
    pdf_stream = io.BytesIO(pdf_bytes)
    pdf_reader = PdfReader(pdf_stream)
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def get_docx_text(docx_bytes):
    text = ""
    doc_stream = io.BytesIO(docx_bytes)
    doc = Document(doc_stream)
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def get_transcript_from_captions(url):
    try:
        video_id = url.split('=')[1]
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        text = ""
        for line in transcript:
            text += " " + line["text"] + "\n"
        return text
    except Exception as e:
        st.error(f"Error retrieving transcript: {str(e)}")
        return None

def process_google_slides(slide_data):
    try:
        prs = Presentation(io.BytesIO(slide_data.read()))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"Error processing Google Slide data: {str(e)}")
        return None

def get_text_chunks(text):
    if not text:
        return []
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=10000, chunk_overlap=1000
    )
    chunks = text_splitter.split_text(text)
    return chunks

def store_text_in_qdrant(text_chunks):
    try:
        embeddings = OpenAIEmbeddings(
            model="text-embedding-ada-002", openai_api_key=OPENAI_API_KEY
        )
        vector_store.add_texts(text_chunks, embeddings=embeddings)
        st.success("Data stored in Qdrant successfully.")
    except Exception as e:
        st.error(f"Error storing data in Qdrant: {str(e)}")

def get_conversational_chain():
    try:
        prompt_template = """
        Answer the question as comprehensively as possible based on the provided context. Ensure to include a summary and all pertinent details related to the question. If the answer cannot be derived from the provided context, simply state, "The answer is not in the context," but refrain from providing incorrect information.\n\n
        Context:\n {context}?\n
        Question:\n {question}\n
        Answer:
        """
        model = ChatOpenAI(model='gpt-3.5-turbo',
                           temperature=0.3, openai_api_key=OPENAI_API_KEY)
        prompt = PromptTemplate(template=prompt_template, input_variables=[
                                "context", 'question'])
        chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
        return chain
    except Exception as e:
        st.error(f"Error loading conversational chain: {str(e)}")
        return None

def process_input(input_data, input_type):
    raw_text = None

    if input_type == "pdf":
        raw_text = get_pdf_text(input_data.getvalue())
    elif input_type == "docx":
        raw_text = get_docx_text(input_data.getvalue())
    elif input_type == "youtube":
        raw_text = get_transcript_from_captions(input_data)
    elif input_type == "ppt":
        raw_text = process_google_slides(input_data)

    if raw_text:
        text_chunks = get_text_chunks(raw_text)
        store_text_in_qdrant(text_chunks)
        st.success(f"Processed {input_type} successfully.")
    else:
        st.error(f"Could not process the {input_type} input.")

def detect_input_type(input_data):
    if isinstance(input_data, str):
        if re.match(r"(https?://)?(www\.)?(youtube\.com|youtu\.be)", input_data):
            return "youtube"
    elif input_data.name.endswith(".pdf"):
        return "pdf"
    elif input_data.name.endswith(".docx"):
        return "docx"
    elif input_data.name.endswith((".pptx", ".ppt")):
        return "ppt"
    return None

def user_input(user_question):
    try:
        docs = vector_store.similarity_search(user_question)

        chain = get_conversational_chain()
        if chain:
            response = chain(
                {"input_documents": docs, "question": user_question}, return_only_outputs=True)
            st.write("Reply: ", response["output_text"])
        else:
            st.error("Error processing user input.")
    except Exception as e:
        st.error(f"Error processing user input: {str(e)}")

def main():
    try:
        st.set_page_config(page_title="MyThoughts.AI",
                           layout='wide', initial_sidebar_state='auto')
        st.title(
            "Welcome to NIFO")
        st.write("Interact with your data")

        st.write("Upload a file (PDF, DOCX, PPT) or enter a YouTube link:")
        input_data = st.file_uploader(
            "File Input", type=["pdf", "docx", "pptx", "ppt"])
         
        user_text = st.text_input("Or enter a YouTube link")

        if st.button("Submit & Process"):
            with st.spinner("Processing..."):
                input_type = detect_input_type(user_text or input_data)
                if input_type:
                    process_input(
                        input_data if input_data else user_text, input_type)
                else:
                    st.error(
                        "Invalid input. Please upload a supported file or provide a valid YouTube link.")

        user_question = st.text_input(
            "Ask a Question about the processed data")

        if user_question:
            user_input(user_question)
    except Exception as e:
        st.error(f"An error occurred: {str(e)}")

if __name__ == "__main__":
    main()
